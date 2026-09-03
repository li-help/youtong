# -*- coding: utf-8 -*-
# QA: 检查文本溢出、越界、文本框重叠
import math, sys
from pptx import Presentation
from pptx.util import Emu

PATH = sys.argv[1] if len(sys.argv) > 1 else "优童成长社-项目答辩PPT.pptx"
EMU = 914400.0
prs = Presentation(PATH)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst) if False else len(prs.slides._sldIdLst))

def char_w(ch, fs, bold):
    if ord(ch) > 0x2E80:
        w = fs / 72.0
    elif ch in "　":
        w = fs / 72.0
    else:
        w = fs * 0.52 / 72.0
    return w * (1.06 if bold else 1.0)

issues = []
for si, slide in enumerate(prs.slides, 1):
    boxes = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        x, y = sh.left / EMU, sh.top / EMU
        w, h = sh.width / EMU, sh.height / EMU
        if x < -0.01 or y < -0.01 or x + w > 13.35 or y + h > 7.52:
            issues.append(f"S{si} OUT-OF-BOUNDS ({x:.2f},{y:.2f},{w:.2f},{h:.2f}) {txt[:18]!r}")
        # 估算所需高度
        need_h = 0.0
        for p in sh.text_frame.paragraphs:
            ptxt = "".join(r.text for r in p.runs)
            if not ptxt.strip():
                need_h += 0.08
                continue
            fs = 18.0
            bold = False
            for r in p.runs:
                if r.font.size:
                    fs = r.font.size.pt
                if r.font.bold:
                    bold = True
                break
            total = sum(char_w(c, fs, bold) for c in ptxt)
            usable = max(w - 0.06, 0.3)
            lines = max(1, math.ceil(total / usable))
            lh = fs * 1.32 / 72.0
            sa = 0.0
            if p.space_after:
                sa = p.space_after.pt / 72.0
            need_h += lines * lh + sa
        if need_h > h * 1.15 + 0.08:
            issues.append(f"S{si} OVERFLOW need={need_h:.2f} box_h={h:.2f} ({x:.2f},{y:.2f},{w:.2f}) {txt[:22]!r}")
        boxes.append((x, y, w, h, txt[:14]))

    # 文本-文本重叠
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ix = max(0, min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0]))
            iy = max(0, min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1]))
            if ix > 0.08 and iy > 0.08:
                issues.append(f"S{si} TEXT-OVERLAP {a[4]!r} x {b[4]!r} area={ix*iy:.2f}")

print(f"\n{len(issues)} issues")
for i in issues:
    print(" ", i)
